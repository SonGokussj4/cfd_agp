"""Evektor library for all things"""
import collections
import csv
import glob
import os
import sys
from collections import namedtuple

import colorlog
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import pptx
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

formatter = colorlog.ColoredFormatter(
    "%(log_color)s%(levelname)-8s%(reset)s %(message_log_color)s%(message)s",
    datefmt=None,
    reset=True,
    log_colors={
        "CRITICAL": "bold_red",
        "ERROR": "red",
        "WARNING": "yellow",
        "INFO": "green",
        "DEBUG": "cyan",
    },
    secondary_log_colors={
        "message": {
            "CRITICAL": "bold_red",
            "ERROR": "red",
            "WARNING": "yellow",
            "INFO": "white",
            "DEBUG": "cyan",
        }
    },
    # style='%',
)

# Initialize LOGGER
handler = colorlog.StreamHandler()
handler.setFormatter(formatter)
logger = colorlog.getLogger(__name__)
logger.addHandler(handler)
logger.setLevel("INFO")


class Line:
    def __init__(self, x1, y1, x2, y2):
        self.x1 = x1
        self.y1 = y1
        self.x2 = x2
        self.y2 = y2
        self.m, self.b = np.polyfit([x1, x2], [y1, y2], 1)

    def get_z(self, coord_x):
        x = float(coord_x)
        return self.m * x + self.b


class Presentation:
    def __init__(self, src_prs_path: str):
        logger.info("Loading template: {}".format(src_prs_path))
        self.prs = pptx.Presentation(src_prs_path)
        self.slides = []
        self.variants = []
        self.conf = None
        self.one_image_slides = [2, 3]
        self.two_images_slides = [4, 5, 6, 7, 8, 9, 10, 11]

    def load_config(self, config_file):
        logger.info("Loading slides configuration: {}".format(config_file))
        import configparser

        conf = configparser.ConfigParser()
        conf.read(config_file)
        self.conf = conf

    def add_variants(self, variants: list):
        if len(variants) not in range(1, 4):
            raise Exception("Only 1-3 variants possible for now...\n")

        for var in variants:
            # Create a NamedTuple
            variant = namedtuple("variant", ("name", "fullpath", "num", "exists"))

            # Check if we have simple path string or dictionary
            if isinstance(var, str):
                var_label = ""
                var_path = var
            else:
                var_label = var.get("label", "")
                var_path = var.get("path")

            # Fill the NamedTuple
            variant.name = var_path.rstrip("/")
            variant.fullpath = os.path.abspath(var_path)
            variant.folder = "/".join(variant.fullpath.split("/")[0:-1])
            if not var_label:
                var_label = variant.name.split("-")[0]
            variant.num = var_label
            variant.exists = os.path.isdir(os.path.abspath(var_path))

            # Check if variant (folder) exists, if true, add to list of Presentation.variants
            if variant.exists:
                self.variants.append(variant)
                logger.info("Added variant {0} in {1}".format(var_label, var_path))
            else:
                logger.critical(
                    "Project: {} was not found in: {} \nAre you in the RIGHT folder??? "
                    "And please check entered variant names. ".format(variant.name, os.getcwd())
                )
                sys.exit()

    def get_num_of_slides(self) -> int:
        return len(self.prs.slides)

    def process_slides(self):
        author = self.conf.get("User Settings", "author")

        for section in self.conf.sections():  # nebo: self.conf.sections()[1:]
            if not self.conf.has_option(section, "layout"):
                continue

            # Convert argparse list of tuples to dictionary of key: val
            conf = dict(self.conf.items(section))

            # Get variables from config file of actual section
            title = conf.get("title", "{} TITLE MISSING".format(section))
            layout_num = int(conf.get("layout", 2))
            fringebar = conf.get("fringebar", None)
            images = conf.get("images", []).replace(" ", "").split(",")

            # Add slide as object
            slide_layout = self.prs.slide_layouts[layout_num]
            pptx_slide = self.prs.slides.add_slide(slide_layout)

            self.slides.append(pptx_slide)
            logger.info("Adding slide {}".format(self.get_num_of_slides()))
            logger.debug(
                "Variables: \nTitle: {} \nLayout: {} \nFringebar: {} \nImages: {}".format(
                    title, layout_num, fringebar, images
                )
            )

            # Slide object: add title, fringebar and images
            slide = Slide(self, pptx_slide, layout_num)
            slide.set_title(title)
            slide.set_author(author)
            slide.add_fringebar(fringebar)
            slide.add_images(images)

    def save_presentation(self, output_pres_path):
        if "/" in output_pres_path:
            output_path = output_pres_path
            self.prs.save(output_path)

        else:
            cur_dir = os.path.realpath(os.path.curdir)
            output_path = os.path.join(cur_dir, output_pres_path)
            self.prs.save(output_path)

        logger.info("Presentation saved to: {}".format(output_path))

    def output_placeholders_pptx(self, output_pres_path: str):
        for master_slide in self.prs.slide_masters:
            for lay_id, layout in enumerate(master_slide.slide_layouts):
                # if lay_id == 0:
                #     continue
                slide = self.prs.slides.add_slide(layout)

                for ph in slide.placeholders:
                    # TEXT modification
                    p = ph.text_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    run = p.add_run()
                    run.font.size = Pt(11)
                    run.font.bold = True

                    # COLOR modification
                    ph.fill.solid()
                    ph.fill.fore_color.rgb = RGBColor(241, 241, 241)

                    if ph.placeholder_format.type == 1:  # TITLE
                        run.text = "Id[{}]: {} - Layout[{}]: [{}]".format(
                            ph.placeholder_format.idx, ph.name, lay_id, layout.name
                        )
                    else:
                        run.text = "Id[{}]: {}".format(ph.placeholder_format.idx, ph.name)

        self.prs.save(output_pres_path)

    def gradients_from_file(self, grad_file):
        os.chdir("PICTURES")
        grad_file_base = "_".join(grad_file.split("_")[0:-1])
        files = glob.glob("{}_[0-9][.][0-9][0-9][0-9]".format(grad_file_base))
        line_win = Line(0.655, 0.688, 0.8, 0.696)
        line_door = Line(0.655, 0.67, 0.8, 0.678)

        # NamedTuple for better manipulation with data
        linetuple = collections.namedtuple("line", ["idx", "dist", "val", "zcoord"])

        win_dict = {}
        door_dict = {}

        for file in files:
            # file = '/ST/SkodaAuto/AEROAKUSTIKA/PRJ/SK370-3/STACIONARNI-VYPOCET/S200-BASIC-MIRROR-DDKM1/PICTURES/Ux_GRAD_0.800'
            x_coord = file.split("_")[-1]
            logger.info("Reading file: {}".format(file))

            # Load data into Pandas DataFrame
            df = pd.read_table(file, comment="$", delimiter=",", engine="python", skipfooter=1, names=["val", "z"])
            df_group = df.groupby("z", sort=False).apply(lambda x: tuple(x["val"])).reset_index()
            res = zip(df_group.get("z").tolist(), df_group.get(0).tolist())

            # Create a NamedTuple from loaded data
            data = [linetuple(idx=idx, val=val, dist=vals[0], zcoord=vals[1]) for idx, (val, vals) in enumerate(res)]

            for idx, line in enumerate(data):
                # Ignore first and last
                if idx == 0 or idx == len(data) - 1:
                    continue

                prev_num = data[idx - 1].val
                y_num = line.val
                next_num = data[idx + 1].val

                # Find where it intersect with 0 y axis
                if prev_num > 0 and y_num < 0 or prev_num < 0 and y_num > 0:
                    ls = [prev_num, y_num, next_num]

                    # Find the minimum absolute number of found elements
                    y_num = min([abs(x) for x in ls])
                    res = [line.idx for line in data if abs(line.val) == abs(y_num) if line.idx != 0]

                    # Edge conditional (it's likely on the edge of line)
                    if res == []:
                        continue

                    new_idx = res[0]

                    # Create list of 3 points of lines by updated new_idx, one of each side of the middle y_num
                    x_coords = [line.dist for line in data[new_idx - 1 : new_idx + 2]]
                    y_coords = [line.val for line in data[new_idx - 1 : new_idx + 2]]

                    if line.zcoord < line_win.get_z(x_coord) and line.zcoord > line_door.get_z(x_coord):
                        location = "Lista"
                        m, b = 0, 0

                    elif line.zcoord >= line_win.get_z(x_coord):
                        location = "Okno"
                        m, b = np.polyfit(x=x_coords, y=y_coords, deg=1)
                        m = abs(m)  # Absolute value
                        if not win_dict.get(x_coord, False):
                            win_dict[x_coord] = m
                        else:
                            if abs(win_dict[x_coord]) < abs(m):
                                win_dict[x_coord] = m

                    elif line.zcoord <= line_door.get_z(x_coord):
                        location = "Dvere"
                        m, b = np.polyfit(x=x_coords, y=y_coords, deg=1)
                        m = abs(m)  # Absolute value
                        if not door_dict.get(x_coord, False):
                            door_dict[x_coord] = m
                        else:
                            if abs(door_dict[x_coord]) < abs(m):
                                door_dict[x_coord] = m

                    else:
                        location = "NEZNAMA CHYBA"

                    logger.info(
                        "{loc:<5} Intersection: [{prev:>12} > {cur:>12} < {next:>12}] {polyfit}".format(
                            loc=location,
                            prev=data[new_idx - 1].val,
                            cur=data[new_idx].val,
                            next=data[new_idx + 1].val,
                            polyfit="... Polyfit: m: {m:.0f}, b: {b:.0f}".format(m=m, b=b) if m != 0 else " ",
                        )
                    )

        win_sorted = sorted(win_dict.items())
        door_sorted = sorted(door_dict.items())

        if win_sorted != []:
            wx, yy = zip(*win_sorted)
            plt.plot(wx, yy, label="Window")

        with open("Ux_GRAD_results_WINDOW", "w") as f:
            wr = csv.writer(f, quoting=csv.QUOTE_NONE)
            # wr.writerows(win_sorted)
            logger.info("WINDOW data saved to: \n{}".format(os.path.abspath(f.name)))

        if door_sorted != []:
            dx, dy = zip(*door_sorted)
            plt.plot(dx, dy, label="Door")

        with open("Ux_GRAD_results_DOOR", "w") as f:
            wr = csv.writer(f, quoting=csv.QUOTE_NONE)
            # wr.writerows(door_sorted)
            logger.info("DOOR data saved to: \n{}".format(os.path.abspath(f.name)))

        plt.legend(loc="upper left", frameon=True)
        plt.xlabel("X_Coordinate")
        plt.ylabel("gradUx(m/s)")

        plt.show()

    def plot_gradients(self):
        plt_colors = ("blue", "red", "violet")
        linreg_rozpeti = 1

        for sec in self.conf.options("Plots"):  # each plot
            sec_name = self.conf.get("Plots", sec)
            fig = plt.figure()
            plt.clf()
            fig.clf()
            plt.style.use("seaborn-notebook")
            plt.grid(True)
            plt.title(sec_name)
            for var_idx, variant in enumerate(self.variants):
                datafile = os.path.join(variant.fullpath, "PICTURES", sec_name)
                with open(datafile, "r") as f:
                    data = f.readlines()

                x_axis = ""
                y_axis = ""
                x_data = []
                y_data = []

                for line in data:
                    line = line.rstrip("\n")
                    if "(X axis) " in line:
                        x_axis = line.split("(X axis) ")[-1]
                    if "(Y axis) " in line:
                        y_axis = line.split("(Y axis) ")[-1]
                    if line.startswith(" "):
                        x_data.append(float(line.replace(" ", "").split(",")[0]))
                        y_data.append(float(line.replace(" ", "").split(",")[1]))

                plt.plot(x_data, y_data, color=plt_colors[var_idx], label=variant.name, linewidth=2.5)
                plt.legend(loc="upper left", frameon=True)
                plt.xlabel(x_axis)
                plt.ylabel(y_axis)
                axes = plt.gca()
                axes.set_ylim([plt.ylim()[0], plt.ylim()[1] + 5])
                axes.invert_xaxis()

                print("\nSec: {} / Variant: {}".format(sec_name, variant.name))

                for idx, num in enumerate(y_data):
                    if idx in (0, len(y_data) - 1):
                        continue

                    prev_num = y_data[idx - 1]
                    next_num = y_data[idx + 1]

                    if prev_num > 0 and num < 0 or prev_num < 0 and num > 0:
                        ls = [prev_num, num, next_num]
                        num = min([abs(x) for x in ls])
                        try:
                            new_index = y_data.index(num)
                        except ValueError:
                            new_index = y_data.index(-num)
                        idx = new_index

                        x = x_data[idx - linreg_rozpeti : idx + linreg_rozpeti + 1]
                        y = y_data[idx - linreg_rozpeti : idx + linreg_rozpeti + 1]
                        m, b = np.polyfit(x, y, 1)
                        print(
                            "Protnuti X: [{} > {} < {}] ... Polyfit: m: {:.0f}, b: {:.0f}".format(
                                prev_num, num, next_num, m, b
                            )
                        )
                        plt.annotate(
                            "{:.0f}".format(m),
                            xy=(x_data[idx], 0),
                            xycoords="data",
                            color=plt_colors[var_idx],
                            xytext=(+15, +15 + var_idx * 15),
                            textcoords="offset points",
                            fontsize=10,
                            bbox=dict(facecolor="white", edgecolor="None", alpha=0.65),
                            arrowprops=dict(arrowstyle="->", connectionstyle="arc3,rad=.2"),
                        )
                    else:
                        continue
                fig.savefig("{}.png".format(sec_name), dpi=800, bbox_inches="tight")

            # plt.show()


class Slide:
    def __init__(self, pres, slide, layout_num):
        self.slide = slide  # slide knows about pptx.slide object
        self.pres = pres  # slide knows about presentation
        self.variants = pres.variants
        self.layout_num = layout_num
        self.slide_num = pres.get_num_of_slides()

    def set_title(self, title: str):
        try:
            self.slide.shapes.title.text = title
        except AttributeError:
            pass

    def set_author(self, author: str):
        try:
            self.slide.placeholders[20].text = author
        except KeyError:
            pass

    def add_images(self, images):
        # Should be 1 image but more were specified in config file
        if self.layout_num in self.pres.one_image_slides and len(images) > 1:
            logger.critical(
                "You've specified [{} images] for layout[{}]. Should be [{} image]. "
                "Fix the config file.".format(len(images), self.layout_num, 1)
            )
            sys.exit()

        # Should be more images but less was specivied in config file
        elif self.layout_num in self.pres.two_images_slides and len(images) < 2:
            logger.critical(
                "You've specified only [{} image] for layout[{}]. Should be [{} images]. "
                "Fix the config file.".format(len(images), self.layout_num, 2)
            )
            sys.exit()

        for idx, variant in enumerate(self.variants):
            # TEXT
            try:
                self.slide.placeholders[17 + idx].text = variant.num  # Variant number
            except KeyError:
                pass

            # IMAGES
            # 1st image is in all slide layouts
            if self.layout_num in self.pres.one_image_slides or self.layout_num in self.pres.two_images_slides:
                img1_path = os.path.join(variant.fullpath, "PICTURES", images[0])
                if os.path.isfile(img1_path):
                    self.slide.placeholders[11 + idx].insert_picture(img1_path)
                else:
                    logger.error(
                        "Image: {} does not exist in \n         {}".format(
                            os.path.basename(img1_path), os.path.dirname(img1_path)
                        )
                    )

            # 2nd additional image is only in layout 2, 4, 5, 8
            if self.layout_num in self.pres.two_images_slides:
                img2_path = os.path.join(variant.fullpath, "PICTURES", images[1])

                if os.path.isfile(img2_path):
                    self.slide.placeholders[14 + idx].insert_picture(img2_path)
                else:
                    logger.error(
                        "Image: {} does not exist in \n         {}".format(
                            os.path.basename(img2_path), os.path.dirname(img2_path)
                        )
                    )

            # if self.layout_num == 1:
            #     img_orig_path = os.path.join(variant.fullpath, 'PICTURES', images[0])
            #     if os.path.isfile(img_orig_path):
            #         self.slide.placeholders[10].insert_picture(img_orig_path)
            #     else:
            #         logger.error("Image: {} does not exist in \n         {}".format(
            #             os.path.basename(img_orig_path), os.path.dirname(img_orig_path)))

            # 1 image on whole slide in layout 12
            if self.layout_num == 12:
                img_orig_path = os.path.join(variant.fullpath, "PICTURES", images[0])
                if os.path.isfile(img_orig_path):
                    self.slide.placeholders[11].insert_picture(img_orig_path, crop=False)
                else:
                    logger.error(
                        "Image: {} does not exist in \n         {}".format(
                            os.path.basename(img_orig_path), os.path.dirname(img_orig_path)
                        )
                    )

            # 6 original images in layout 13
            if self.layout_num == 13:
                for NUM in range(0, 6):
                    img_orig_path = os.path.join(variant.fullpath, "PICTURES", images[NUM])
                    if os.path.isfile(img_orig_path):
                        self.slide.placeholders[11 + NUM].insert_picture(img_orig_path, crop=False)
                    else:
                        logger.error(
                            "Image: {} does not exist in \n         {}".format(
                                os.path.basename(img_orig_path), os.path.dirname(img_orig_path)
                            )
                        )

    def add_fringebar(self, fringebar: str):
        if self.layout_num in [2, 4, 10, 11] and fringebar is None:
            logger.critical(
                "Slide [{}] with Layout[{}] has to have fringebar but none was specified "
                "in config file. Aborting script...".format(self.slide_num, self.layout_num)
            )
            sys.exit()

        elif self.layout_num not in [2, 4, 10, 11] and fringebar is not None:
            logger.error(
                "Slide [{}] with Layout[{}] should not have FRINGEBAR assigned. "
                "Please see the config file and check this slide.".format(self.slide_num, self.layout_num)
            )
            return None

        if fringebar is not None:
            fringebar_path = os.path.join(self.variants[0].fullpath, "PICTURES", fringebar)
            if os.path.isfile(fringebar_path):
                self.slide.placeholders[10].insert_picture(fringebar_path)
            else:
                logger.warning(
                    "Fringebar: {} does not exist in \n         {}".format(
                        os.path.basename(fringebar_path), os.path.dirname(fringebar_path)
                    )
                )
